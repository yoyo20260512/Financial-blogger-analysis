"""
文件解析器 - 支持 PDF、PPT、TXT/MD 文字提取
"""
import io
import logging

logger = logging.getLogger(__name__)


def parse_text(filename: str, content: bytes) -> str:
    """根据文件后缀提取文字内容"""
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext == "txt" or ext == "md":
        return content.decode("utf-8", errors="replace")

    elif ext == "pdf":
        return _parse_pdf(content)

    elif ext in ("ppt", "pptx"):
        return _parse_pptx(content)

    else:
        logger.warning("Unsupported file type: %s", ext)
        return ""


def _parse_pdf(content: bytes) -> str:
    """用 pypdf 提取 PDF 文字"""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        texts = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
        return "\n\n".join(texts)
    except Exception as e:
        logger.warning("PDF parse error: %s", e)
        return ""


def _parse_pptx(content: bytes) -> str:
    """用 python-pptx 提取 PPT 文字"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append("\n".join(slide_texts))
        return "\n\n---\n\n".join(texts)
    except Exception as e:
        logger.warning("PPT parse error: %s", e)
        return ""
